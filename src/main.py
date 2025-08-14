"""
HTML转PowerPoint演示文稿转换器

该脚本将HTML文件转换为PowerPoint演示文稿，支持：
- 多线程并行处理
- 自动提取页面元素（文本、图标、背景）
- 保持原始布局和样式
- 高分辨率图标截图
"""

import os
import time
import concurrent.futures
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import logging
import argparse
import shutil
import re

# ========== 数据结构定义 ==========

class ElementData:
    """
    HTML元素数据容器
    
    存储从HTML页面提取的元素信息，包括：
    - 基本属性（标签名、CSS类）
    - 几何信息（位置、尺寸、样式）
    - 内容（文本、图标路径）
    - 背景信息（是否有背景色）
    - 层级关系（子元素）
    """
    def __init__(self):
        self.tag_name = None                    # HTML标签名
        self.classes = []                       # CSS类名列表
        self.text = None                        # 元素文本内容
        self.geom = None                        # 几何和样式信息
        self.icon_path = None                   # 图标截图路径
        self.has_background = False             # 是否有有效背景色
        self.children = []                      # 子元素列表

class SlideData:
    """
    幻灯片数据容器
    
    存储单个幻灯片的所有信息：
    - 所有页面元素数据
    """
    def __init__(self):
        self.elements = []                      # 页面元素列表

# ========== 配置常量 ==========

# 图标CSS类名集合 - 用于识别需要特殊处理的图标元素
ICON_CLASSES = {
    'material-icons', 'toc-icon', 'importance-icon', 'limitation-icon', 
    'check-icon', 'partial-icon', 'close-icon', 'feature-icon', 'section-icon', 
    'api-icon', 'config-icon', 'case-icon', 'component-icon', 'mock-icon', 
    'snapshot-icon', 'resource-icon'
}

# 布局配置常量
SLIDE_WIDTH_PX = 1280                       # 幻灯片宽度（像素）


# ========== 核心解析逻辑 ==========

def parse_color(color_str):
    """
    解析CSS颜色字符串
    
    Args:
        color_str: CSS颜色字符串，支持rgb()和rgba()格式
        
    Returns:
        tuple: (r, g, b, a) 颜色值元组，解析失败时返回透明黑色(0,0,0,0)
    """
    if not color_str or color_str in ['transparent', 'inherit', 'initial', 'unset']:
        return (0, 0, 0, 0)
    try:
        # 提取颜色值中的数字部分（包括小数）
        parts = re.findall(r'[\d.]+', color_str)
        if len(parts) < 3:
            return (0, 0, 0, 0)
            
        r, g, b = int(float(parts[0])), int(float(parts[1])), int(float(parts[2]))
        a = float(parts[3]) if len(parts) > 3 else 1.0  # alpha通道，默认为1.0（不透明）
        
        # 确保颜色值在有效范围内
        r = max(0, min(255, r))
        g = max(0, min(255, g))
        b = max(0, min(255, b))
        a = max(0.0, min(1.0, a))
        
        return (r, g, b, a)
    except (IndexError, ValueError, TypeError):
        return (0, 0, 0, 0)  # 解析失败时返回透明黑色

def init_driver():
    """
    初始化Selenium WebDriver
    
    配置Chrome浏览器选项：
    - 固定窗口尺寸（1280x720）
    - 隐藏滚动条
    - 无头模式运行（提高性能）
    - 直接使用本地ChromeDriver路径，跳过版本检查
    
    Returns:
        webdriver.Chrome: 配置好的Chrome WebDriver实例
    """
    options = webdriver.ChromeOptions()
    options.add_argument("--window-size=1280,720")  # 设置浏览器窗口尺寸
    options.add_argument("--hide-scrollbars")       # 隐藏滚动条
    options.add_argument('--headless')              # 无头模式，不显示浏览器界面
    
    # 性能优化配置
    options.add_argument("--disable-web-security")  # 禁用网络安全检查
    options.add_argument("--disable-features=VizDisplayCompositor")  # 提高截图稳定性
    options.add_argument("--no-sandbox")            # 提高兼容性
    options.add_argument("--disable-dev-shm-usage") # 避免内存问题
    
    # 直接使用本地ChromeDriver路径，跳过版本检查和下载
    chromedriver_path = r"C:\Users\Administrator\.wdm\drivers\chromedriver\win64\138.0.7204.183\chromedriver-win32\chromedriver.exe"
    
    service = Service(chromedriver_path)
    driver = webdriver.Chrome(service=service, options=options)
    
    logging.info("WebDriver已初始化（无头模式，使用本地ChromeDriver）")
    return driver

def take_icon_screenshot(driver, icon_element, temp_dir, slide_index, element_index, slide_element=None):
    """
    拍摄图标元素的高分辨率截图
    
    通过创建放大的克隆元素来获取清晰的图标截图，并自动裁剪空白区域。
    
    Args:
        driver: Selenium WebDriver实例
        icon_element: 要截图的图标元素
        temp_dir: 临时文件存储目录
        slide_index: 幻灯片索引
        element_index: 元素索引
        slide_element: 父级幻灯片元素（用于避免干扰）
        
    Returns:
        str: 截图文件路径，失败时返回None
    """
    # 临时移动幻灯片元素，避免其他内容干扰截图
    if slide_element:
        try:
            driver.execute_script("arguments[0].style.transform = 'translateY(1000px)';", slide_element)
            time.sleep(0.1)  # 等待DOM更新完成
        except Exception as e:
            logging.warning(f"无法移动slide元素: {e}")
    
    # 创建5倍放大的图标克隆用于高清截图
    scale_factor = 5
    js_script = """
        const targetElement = arguments[0];
        const scale = arguments[1];
        const style = window.getComputedStyle(targetElement);
        const originalFontSizeStr = style.getPropertyValue('font-size');
        const originalColor = style.getPropertyValue('color');
        if (!originalFontSizeStr) return null;
        const originalFontSize = parseFloat(originalFontSizeStr);
        
        // 创建临时容器
        const container = document.createElement('div');
        container.id = 'temp-icon-container-for-screenshot';
        container.style.position = 'absolute';
        container.style.left = '0px';
        container.style.top = '0px';
        container.style.zIndex = '9999';
        
        // 创建放大的图标克隆
        const clone = targetElement.cloneNode(true);
        clone.style.fontSize = (originalFontSize * scale) + 'px';
        clone.style.color = originalColor;
        clone.style.backgroundColor = 'transparent';
        container.appendChild(clone);
        document.body.appendChild(container);
        return clone;
    """
    
    scaled_clone_element = driver.execute_script(js_script, icon_element, scale_factor)
    if not scaled_clone_element:
        logging.warning(f"无法为幻灯片{slide_index}的图标创建可缩放克隆")
        # 恢复幻灯片元素位置
        if slide_element:
            try:
                driver.execute_script("arguments[0].style.transform = '';", slide_element)
            except Exception:
                pass
        return None

    time.sleep(0.1)  # 等待渲染完成
    
    icon_path = os.path.join(temp_dir, f"slide_{slide_index}_element_{element_index}_icon.png")
    try:
        # 截图并自动裁剪空白区域
        logging.info(f"正在截取图标: {icon_path}")
        time.sleep(0.2)  # 等待渲染稳定
        scaled_clone_element.screenshot(icon_path)
        with Image.open(icon_path) as img:
            bbox = img.getbbox()  # 获取非透明区域的边界框
            if bbox:
                cropped_img = img.crop(bbox)
                cropped_img.save(icon_path)
        return icon_path
    except Exception as e:
        logging.error(f"图标截图或裁剪失败: {e}")
        return None
    finally:
        # 清理临时DOM元素
        driver.execute_script("document.getElementById('temp-icon-container-for-screenshot').remove();")
        # 恢复幻灯片元素位置
        if slide_element:
            try:
                driver.execute_script("arguments[0].style.transform = '';", slide_element)
            except Exception as e:
                logging.warning(f"无法恢复slide元素位置: {e}")

def take_code_block_screenshot(driver, code_element, temp_dir, slide_index, element_index, slide_element=None):
    """
    拍摄code-block元素的截图
    
    直接对code-block元素进行截图，保持其原始样式和布局。
    
    Args:
        driver: Selenium WebDriver实例
        code_element: 要截图的code-block元素
        temp_dir: 临时文件存储目录
        slide_index: 幻灯片索引
        element_index: 元素索引
        slide_element: 父级幻灯片元素（用于避免干扰）
        
    Returns:
        str: 截图文件路径，失败时返回None
    """
    # 临时移动幻灯片元素，避免其他内容干扰截图
    if slide_element:
        try:
            driver.execute_script("arguments[0].style.transform = 'translateY(1000px)';", slide_element)
            time.sleep(0.1)  # 等待DOM更新完成
        except Exception as e:
            logging.warning(f"无法移动slide元素: {e}")
    
    code_block_path = os.path.join(temp_dir, f"slide_{slide_index}_element_{element_index}_code_block.png")
    try:
        # 确保元素可见并等待渲染稳定
        logging.info(f"正在截取code-block: {code_block_path}")
        time.sleep(0.2)  # 等待渲染稳定
        
        # 直接对code-block元素截图
        code_element.screenshot(code_block_path)
        
        # 可选：裁剪空白区域（如果需要的话）
        with Image.open(code_block_path) as img:
            bbox = img.getbbox()  # 获取非透明区域的边界框
            if bbox:
                cropped_img = img.crop(bbox)
                cropped_img.save(code_block_path)
        
        return code_block_path
    except Exception as e:
        logging.error(f"code-block截图失败: {e}")
        return None
    finally:
        # 恢复幻灯片元素位置
        if slide_element:
            try:
                driver.execute_script("arguments[0].style.transform = '';", slide_element)
            except Exception as e:
                logging.warning(f"无法恢复slide元素位置: {e}")

def parse_element_recursively(driver, element, temp_dir, slide_index, element_counter, parent_bg_color=None, slide_element=None):
    """
    递归解析HTML元素及其子元素
    
    提取元素的所有相关信息：基本属性、几何信息、文本内容等。
    对于图标元素，会拍摄高分辨率截图；对于有背景色的元素，直接记录背景信息。
    
    Args:
        driver: Selenium WebDriver实例
        element: 要解析的HTML元素
        temp_dir: 临时文件存储目录
        slide_index: 幻灯片索引
        element_counter: 元素计数器（字典，用于跨递归调用保持状态）
        parent_bg_color: 父元素背景色（用于判断是否需要新背景）
        slide_element: 根级幻灯片元素
        
    Returns:
        ElementData: 解析后的元素数据，解析失败时返回None
    """
    data = ElementData()
    
    # 获取元素基本属性
    try:
        data.tag_name = element.tag_name
        data.classes = element.get_attribute('class').split() if element.get_attribute('class') else []
        print(f"🔍 解析元素: <{data.tag_name}> 类名: {data.classes}")
    except Exception:
        print("❌ 无法获取元素基本属性")
        return None

    # 获取元素几何信息和样式属性
    try:
        location = element.location
        size = element.size
        # 跳过不可见元素
        if size['width'] == 0 or size['height'] == 0:
             print(f"⚠️  跳过不可见元素: 尺寸 {size['width']}x{size['height']}")
             return None
             
        data.geom = {
            "x": location['x'], 
            "y": location['y'],
            "width": size['width'], 
            "height": size['height'],
            "font-size": element.value_of_css_property('font-size'),
            "color": element.value_of_css_property('color'),
            "font-weight": element.value_of_css_property('font-weight'),
            "text-align": element.value_of_css_property('text-align'),
            'background-color': element.value_of_css_property('background-color'),
            'border-radius': element.value_of_css_property('border-radius'),
            'box-shadow': element.value_of_css_property('box-shadow')
        }
        print(f"📐 几何信息: 位置({location['x']}, {location['y']}) 尺寸({size['width']}x{size['height']})")
        print(f"🎨 样式信息: 字体大小={data.geom['font-size']} 颜色={data.geom['color']} 背景色={data.geom['background-color']}")
    except Exception:
        print("❌ 无法获取元素几何信息")
        return None  # 元素不可见或无法交互

    # 图标元素特殊处理：拍摄高分辨率截图后直接返回（图标是叶子节点）
    if any(cls in ICON_CLASSES for cls in data.classes):
        print(f"🎯 发现图标元素: {data.classes}")
        data.icon_path = take_icon_screenshot(driver, element, temp_dir, slide_index, element_counter['i'], slide_element)
        if data.icon_path:
            print(f"📸 图标截图成功: {data.icon_path}")
        else:
            print("❌ 图标截图失败")
        return data

    # code-block元素特殊处理：直接截图，不解析内部结构
    if 'code-block' in data.classes:
        print(f"💻 发现code-block元素: {data.classes}")
        data.icon_path = take_code_block_screenshot(driver, element, temp_dir, slide_index, element_counter['i'], slide_element)
        if data.icon_path:
            print(f"📸 code-block截图成功: {data.icon_path}")
        else:
            print("❌ code-block截图失败")
        return data

    # 提取文本内容
    try:
        # 使用JavaScript提取直接文本节点内容
        js_get_text = """
        function calculateWidth(text) {
            let width = 0;
            for (let char of text) {
            // 中文字符算作2个空格宽度，英文字符算作1个
            width += /[\\u4e00-\\u9fa5]/.test(char) ? 2 : 1;
            }
            return width;
        }
        
        return Array.from(arguments[0].childNodes)
            .map(node => {
            if (node.nodeType === 3) {
                return node.nodeValue.trim();
            } else if (node.nodeType === 1 && node.tagName === 'SPAN') {
                const spanText = node.innerText.trim();
                const spaceCount = calculateWidth(spanText);
                return ' '.repeat(spaceCount * 2);
            } else if (node.nodeType === 1 && node.tagName === 'BR') {
                return '\\n';
            } else if (node.nodeType === 1 && node.tagName === 'I') {
                return '    ';
            }
            return '';
            })
            .join(' ');
        """
        text = driver.execute_script(js_get_text, element)
        if text:
            data.text = text
            print(f"📝 提取文本: '{text}'")
        else:
            print("📝 无文本内容")
    except Exception as e:
        print(f"❌ 无法提取文本: {e}")
        logging.warning(f"无法从元素提取文本: {e}")

    # 处理背景色信息
    try:
        bg_color_str = data.geom.get('background-color')
        bg_color = parse_color(bg_color_str)
        # 判断是否有有效背景（alpha > 0 即任何透明度都处理，且与父元素背景不同）
        # 对于rgba(211, 47, 47, 0.05)这样的低透明度背景也要正确处理
        data.has_background = bg_color[3] > 0 and bg_color != parent_bg_color
        if data.has_background:
            print(f"🎨 元素有背景色: rgba{bg_color} 位置({data.geom['x']}, {data.geom['y']}) 尺寸({data.geom['width']}x{data.geom['height']})")
            logging.info(f"元素有背景色: rgba{bg_color} 位置({data.geom['x']}, {data.geom['y']}) 尺寸({data.geom['width']}x{data.geom['height']})")
        else:
            print(f"🔍 无有效背景色: {bg_color_str} -> rgba{bg_color}")
    except Exception:
        data.has_background = False
        print("❌ 无法处理背景色信息")

    # 递归解析子元素
    child_elements = element.find_elements(By.XPATH, "./*")
    print(f"👶 发现 {len(child_elements)} 个子元素")
    current_bg_for_children = parse_color(data.geom.get('background-color', '')) if data.has_background else parent_bg_color
    
    for i, child_element in enumerate(child_elements):
        element_counter['i'] += 1
        print(f"  └─ 处理第 {i+1}/{len(child_elements)} 个子元素 (总计第{element_counter['i']}个)")
        child_data = parse_element_recursively(
            driver, child_element, temp_dir, slide_index, element_counter, 
            parent_bg_color=current_bg_for_children, slide_element=slide_element
        )
        if child_data:
            data.children.append(child_data)
            print(f"  ✅ 子元素解析成功")
        else:
            print(f"  ❌ 子元素解析失败或被跳过")

    # 数据剪枝：如果元素没有任何有用内容，则不返回
    if not data.text and not data.icon_path and not data.has_background and not data.children:
        print("🗑️  元素无有用内容，被剪枝")
        return None

    print(f"✅ 元素解析完成: 文本={bool(data.text)} 图标={bool(data.icon_path)} 背景={data.has_background} 子元素={len(data.children)}")
    return data


def wait_for_material_icons(driver, timeout=10):
    """
    等待Material Icons字体加载完成
    
    确保图标字体完全加载后再进行截图，避免图标显示异常。
    
    Args:
        driver: Selenium WebDriver实例
        timeout: 超时时间（秒）
    """
    js_script = """
        return new Promise((resolve) => {
            if (document.fonts && document.fonts.ready) {
                document.fonts.ready.then(() => {
                    // 额外等待确保material-icons完全加载
                    setTimeout(() => resolve(true), 500);
                });
            } else {
                // 降级处理：如果不支持document.fonts API，等待固定时间
                setTimeout(() => resolve(true), 2000);
            }
        });
    """
    
    try:
        driver.set_script_timeout(timeout)
        driver.execute_async_script(js_script)
        logging.info("Material Icons字体加载完成")
    except Exception as e:
        logging.warning(f"等待Material Icons加载时出错: {e}，继续执行")

def extract_data_from_html(driver, file_path, temp_dir):
    """
    从HTML文件提取所有幻灯片的结构化数据
    
    主要处理流程：
    1. 加载HTML文件并等待渲染完成
    2. 等待字体加载完成
    3. 逐个处理每张幻灯片：
       - 递归解析所有页面元素（跳过背景截图，使用白色背景）
    
    Args:
        driver: 预初始化的Selenium WebDriver实例
        file_path: HTML文件路径
        temp_dir: 临时文件存储目录
        
    Returns:
        list: 包含所有幻灯片数据的列表
    """
    # 加载HTML文件
    driver.get(f"file:///{os.path.abspath(file_path)}")
    time.sleep(2)  # 等待页面渲染完成
    
    # 等待Material Icons字体加载完成
    wait_for_material_icons(driver)

    slides_data = []
    slide_elements = driver.find_elements(By.CSS_SELECTOR, ".slide")
    logging.info(f"在文件 {os.path.basename(file_path)} 中找到 {len(slide_elements)} 张幻灯片")

    for i, slide_element in enumerate(slide_elements):
        logging.info(f"正在处理第 {i+1} 张幻灯片...")
        slide_data = SlideData()
        element_counter = {'i': 0}  # 使用可变字典作为计数器，在递归中保持状态

        # 递归解析幻灯片内容（跳过背景截图，因为都是白色背景）
        try:
            # 查找并处理所有主要容器：slide-header 和 slide-content
            containers_to_process = []
            
            # 查找 slide-header 容器
            try:
                header_element = slide_element.find_element(By.CSS_SELECTOR, ".slide-header")
                containers_to_process.append(header_element)
                logging.info(f"第 {i+1} 张幻灯片找到 slide-header 容器")
            except Exception:
                logging.info(f"第 {i+1} 张幻灯片未找到 slide-header 容器")
            
            # 查找 slide-content 容器
            try:
                content_element = slide_element.find_element(By.CSS_SELECTOR, ".slide-content")
                containers_to_process.append(content_element)
                logging.info(f"第 {i+1} 张幻灯片找到 slide-content 容器")
            except Exception:
                logging.info(f"第 {i+1} 张幻灯片未找到 slide-content 容器")
            
            # 如果没有找到任何容器，直接从 slide 根元素解析
            if not containers_to_process:
                logging.warning(f"第 {i+1} 张幻灯片未找到任何标准容器，从 slide 根元素开始解析")
                containers_to_process.append(slide_element)
            
            # 处理所有找到的容器
            for container in containers_to_process:
                child_elements = container.find_elements(By.XPATH, "./*")
                for child in child_elements:
                    element_counter['i'] += 1
                    element_data = parse_element_recursively(
                        driver, child, temp_dir, i, element_counter, slide_element=slide_element
                    )
                    if element_data:
                        slide_data.elements.append(element_data)
                        
        except Exception as e:
            logging.error(f"处理第 {i+1} 张幻灯片内容时出错: {e}")

        slides_data.append(slide_data)

    return slides_data



# ========== PowerPoint生成逻辑 ==========

def create_presentation():
    """
    创建PowerPoint演示文稿
    
    设置标准的16:9宽屏比例（13.333" x 7.5"）
    
    Returns:
        Presentation: 配置好的PowerPoint演示文稿对象
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9比例的宽度
    prs.slide_height = Inches(7.5)     # 16:9比例的高度
    return prs

def px_to_emu(px):
    """
    像素转换为EMU单位
    
    EMU (English Metric Units) 是PowerPoint内部使用的度量单位。
    1像素 = 9525 EMU
    
    Args:
        px: 像素值
        
    Returns:
        int: EMU值
    """
    if isinstance(px, (int, float)):
        return int(px * 9525)
    return 0

def add_slide_with_white_background(prs):
    """
    添加白色背景的幻灯片
    
    Args:
        prs: PowerPoint演示文稿对象
        
    Returns:
        Slide: 创建的幻灯片对象
    """
    slide_layout = prs.slide_layouts[6]  # 使用空白布局（默认白色背景）
    slide = prs.slides.add_slide(slide_layout)
    return slide

def add_image(slide, image_path, geom):
    """
    向幻灯片添加图片元素
    
    根据几何信息精确定位和调整图片大小。
    
    Args:
        slide: 幻灯片对象
        image_path: 图片文件路径
        geom: 几何信息字典（包含位置和尺寸）
    """
    if not (image_path and os.path.exists(image_path) and geom):
        return
        
    try:
        # 转换像素坐标为EMU单位
        x_emu = px_to_emu(geom['x'])
        y_emu = px_to_emu(geom['y'])
        width_emu = px_to_emu(geom['width'])
        height_emu = px_to_emu(geom['height'])
        
        # 确保最小尺寸，避免零尺寸元素导致错误
        if width_emu == 0 or height_emu == 0:
            default_size = px_to_emu(20)  # 微小/不可见图标的默认尺寸
            width_emu = width_emu or default_size
            height_emu = height_emu or default_size

        slide.shapes.add_picture(image_path, x_emu, y_emu, width=width_emu, height=height_emu)
    except Exception as e:
         logging.error(f"添加图片失败 {image_path}: {e}")

def parse_border_radius(border_radius_str):
    """
    解析CSS border-radius值
    
    Args:
        border_radius_str: CSS border-radius字符串，如 "12px" 或 "8px 4px"
        
    Returns:
        float: 圆角半径像素值，解析失败返回0
    """
    if not border_radius_str or border_radius_str in ['0', '0px', 'none']:
        return 0
    try:
        # 提取第一个数字值（简化处理，只取第一个圆角值）
        match = re.search(r'(\d+(?:\.\d+)?)', border_radius_str)
        if match:
            return float(match.group(1))
    except (ValueError, AttributeError):
        pass
    return 0

def add_background_shape(slide, geom):
    """
    向幻灯片添加背景形状，支持圆角和透明度
    
    根据CSS背景色、圆角等信息创建形状，支持低透明度背景色的正确处理。
    
    Args:
        slide: 幻灯片对象
        geom: 几何信息字典（包含位置、尺寸、背景色、圆角等）
    """
    if not geom or 'background-color' not in geom:
        return
        
    try:
        # 解析背景色
        bg_color = parse_color(geom['background-color'])
        if bg_color[3] <= 0:  # 完全透明时跳过
            return
            
        # 转换像素坐标为EMU单位
        x_emu = px_to_emu(geom['x'])
        y_emu = px_to_emu(geom['y'])
        width_emu = px_to_emu(geom['width'])
        height_emu = px_to_emu(geom['height'])
        
        # 确保最小尺寸
        if width_emu <= 0 or height_emu <= 0:
            return
            
        # 解析圆角
        border_radius = parse_border_radius(geom.get('border-radius', '0'))
        
        # 创建形状（根据是否有圆角选择形状类型）
        from pptx.enum.shapes import MSO_SHAPE
        if border_radius > 0:
            # 有圆角时使用圆角矩形
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_emu, y_emu, width_emu, height_emu
            )
            # 设置圆角半径（PowerPoint中的调整值需要转换）
            try:
                # PowerPoint的圆角调整值范围通常是0-1之间
                corner_radius = min(0.5, border_radius / min(geom['width'], geom['height']) * 2)
                shape.adjustments[0] = corner_radius
            except Exception as e:
                logging.warning(f"设置圆角失败: {e}")
        else:
            # 无圆角时使用普通矩形
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x_emu, y_emu, width_emu, height_emu
            )
        
        # 设置填充色 - 对于透明度，直接计算与白色背景的混合色
        fill = shape.fill
        fill.solid()
        
        if bg_color[3] < 1.0:
            # 计算与白色背景的混合色，而不是依赖PowerPoint的透明度功能
            # Alpha混合公式: result = foreground * alpha + background * (1 - alpha)
            # 白色背景 RGB(255, 255, 255)
            alpha = bg_color[3]
            mixed_r = int(bg_color[0] * alpha + 255 * (1 - alpha))
            mixed_g = int(bg_color[1] * alpha + 255 * (1 - alpha))
            mixed_b = int(bg_color[2] * alpha + 255 * (1 - alpha))
            
            fill.fore_color.rgb = RGBColor(mixed_r, mixed_g, mixed_b)
            logging.info(f"原始颜色: rgba{bg_color}, 混合后颜色: RGB({mixed_r}, {mixed_g}, {mixed_b})")
        else:
            # 完全不透明，直接使用原始颜色
            fill.fore_color.rgb = RGBColor(bg_color[0], bg_color[1], bg_color[2])
            logging.info(f"不透明颜色: RGB({bg_color[0]}, {bg_color[1]}, {bg_color[2]})")
        
        # 移除边框
        line = shape.line
        line.fill.background()
        
        # 处理阴影效果（简化处理）
        box_shadow = geom.get('box-shadow', '')
        if box_shadow and box_shadow != 'none':
            try:
                # 简单的阴影效果
                shadow = shape.shadow
                shadow.inherit = False
                shadow.style = 1  # 外阴影
                shadow.blur_radius = px_to_emu(4)  # 默认模糊半径
                shadow.distance = px_to_emu(2)    # 默认阴影距离
                shadow.color.rgb = RGBColor(0, 0, 0)  # 黑色阴影
                shadow.transparency = 0.8  # 阴影透明度
            except Exception as e:
                logging.warning(f"设置阴影失败: {e}")
        
        logging.info(f"添加背景形状: 位置({geom['x']}, {geom['y']}) 尺寸({geom['width']}x{geom['height']}) 颜色rgba{bg_color} 圆角{border_radius}px")
        
    except Exception as e:
        logging.error(f"添加背景形状失败: {e}")

def add_textbox(slide, data, slide_width_px):
    """
    向幻灯片添加文本框，支持动态宽度调整
    
    为所有文本元素添加安全缓冲区防止文本换行
    
    Args:
        slide: 幻灯片对象
        data: 元素数据对象
        slide_width_px: 幻灯片宽度（像素）
    """
    if not data or not data.text or not data.geom:
        return

    geom = data.geom.copy()  # 使用副本避免副作用
    text = data.text

    # ========== 统一宽度调整逻辑 ==========
    # 为所有元素添加安全缓冲区防止文本换行
    geom['width'] += 30

    try:
        # 转换坐标和尺寸
        x_emu = px_to_emu(geom['x'])
        y_emu = px_to_emu(geom['y'])
        width_emu = px_to_emu(geom['width'])
        height_emu = px_to_emu(geom['height'])

        # 创建文本框
        textbox = slide.shapes.add_textbox(x_emu, y_emu, width_emu, height_emu)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text

        # 应用字体样式
        font = run.font
        
        # 设置字体大小
        if 'font-size' in geom:
            try:
                font_size_px = float(str(geom['font-size']).replace("px", ""))
                scale_factor = 0.75  # PowerPoint渲染的调整系数
                font.size = Pt(int(font_size_px * scale_factor))
            except (ValueError, TypeError):
                pass
                
        # 设置字体颜色
        if 'color' in geom:
            try:
                color_str = geom['color'].replace("rgba(", "").replace("rgb(", "").replace(")", "")
                parts = [p.strip() for p in color_str.split(",")]
                r, g, b = int(parts[0]), int(parts[1]), int(parts[2])
                font.color.rgb = RGBColor(r, g, b)
            except Exception:
                pass
                
        # 设置字体粗细
        if 'font-weight' in geom:
            font_weight = str(geom['font-weight'])
            if font_weight == 'bold' or (font_weight.isnumeric() and int(font_weight) >= 700):
                font.bold = True
                
    except Exception as e:
        logging.error(f"添加文本框失败，文本内容: '{text}': {e}")

def add_elements_to_slide(slide, elements, slide_width_px):
    """
    递归地将元素数据添加到幻灯片
    
    处理顺序：
    1. 元素背景形状（如果有）
    2. 元素内容（图标或文本）
    3. 递归处理子元素
    
    Args:
        slide: 幻灯片对象
        elements: 元素数据列表
        slide_width_px: 幻灯片宽度（像素）
    """
    for data in elements:
        # 1. 添加元素背景形状（如果有）
        if data.has_background:
            add_background_shape(slide, data.geom)

        # 2. 添加元素内容（图标或文本）
        if data.icon_path:
            add_image(slide, data.icon_path, data.geom)
        elif data.text:
            add_textbox(slide, data, slide_width_px)

        # 3. 递归处理子元素
        if data.children:
            add_elements_to_slide(slide, data.children, slide_width_px)


# ========== 主执行逻辑 ==========

def process_files_worker(task_info):
    """
    多线程工作函数
    
    每个工作线程初始化一个WebDriver实例，处理分配给它的HTML文件块。
    这样可以避免线程间的WebDriver冲突，提高并行处理效率。
    
    Args:
        task_info: 包含(文件块, 输入目录, 输出目录, 临时目录)的元组
    """
    files_chunk, input_dir, output_dir, temp_dir = task_info
    
    if not files_chunk:
        return  # 如果没有分配文件，直接返回

    logging.info(f"工作线程启动，分配到 {len(files_chunk)} 个文件")
    print(f"\n🚀 开始处理文件: {files_chunk}")
    driver = None
    try:
        driver = init_driver()
        for html_file in files_chunk:
            try:
                print(f"\n📄 开始处理文件: {html_file}")
                logging.info(f"--- 正在处理文件: {html_file} ---")

                # 每个线程使用独立的临时子目录，避免文件冲突
                thread_temp_dir = os.path.join(temp_dir, re.sub(r'[^a-zA-Z0-9.-]', '_', html_file))
                if not os.path.exists(thread_temp_dir):
                    os.makedirs(thread_temp_dir)
                print(f"📁 临时目录: {thread_temp_dir}")

                # 为每个文件创建新的演示文稿
                prs = create_presentation()

                file_path = os.path.join(input_dir, html_file)
                print(f"🔍 开始提取HTML数据: {file_path}")
                # 提取HTML数据
                all_slides_data = extract_data_from_html(driver, file_path, thread_temp_dir)
                print(f"✅ 提取完成，共找到 {len(all_slides_data)} 张幻灯片")

                print(f"🎨 开始生成PowerPoint幻灯片...")
                # 生成PowerPoint幻灯片（使用白色背景）
                for slide_data in all_slides_data:
                    slide = add_slide_with_white_background(prs)
                    add_elements_to_slide(slide, slide_data.elements, SLIDE_WIDTH_PX)

                # 保存演示文稿
                base_name = os.path.splitext(html_file)[0]
                output_path = os.path.join(output_dir, f"{base_name}.pptx")
                print(f"💾 保存PowerPoint文件: {output_path}")
                prs.save(output_path)
                print(f"✅ 成功创建 {output_path}")
                logging.info(f"成功创建 {output_path}")
                
            except Exception as e:
                # 记录单个文件的异常，但继续处理其他文件
                logging.error(f"处理文件 {html_file} 失败: {e}", exc_info=True)
    except Exception as e:
        # 记录整个工作线程的致命异常（如WebDriver初始化失败）
        logging.critical(f"工作线程发生致命错误: {e}", exc_info=True)
    finally:
        if driver:
            print("🔒 关闭浏览器...")
            driver.quit()
            logging.info("工作线程完成，WebDriver已关闭")

def main():
    """
    主函数 - 程序入口点
    
    处理命令行参数，设置工作环境，分配任务给多个工作线程并行处理。
    """
    # 解析命令行参数
    parser = argparse.ArgumentParser(description='将HTML文件转换为PowerPoint演示文稿')
    parser.add_argument('--input_path', type=str, required=True, 
                       help='输入HTML文件路径或包含HTML文件的目录路径')
    parser.add_argument('--output_dir', type=str, required=True, 
                       help='生成的PPTX文件输出目录')
    parser.add_argument('--workers', type=int, default=2, 
                       help='并行转换使用的线程数，默认为2')
    args = parser.parse_args()

    # 配置日志
    logging.basicConfig(
        level=logging.INFO, 
        format='%(asctime)s - %(threadName)s - %(levelname)s - %(message)s'
    )
    
    # 优化模式提示
    print("=" * 60)
    print("🚀 HTML转PowerPoint转换器")
    print("- 使用直接背景形状生成，提高转换效率")
    print("- 保持高质量图标截图")
    print("=" * 60)
    
    input_path = args.input_path
    output_dir = args.output_dir

    # 设置临时目录
    base_dir = os.path.dirname(os.path.abspath(__file__))
    temp_dir = os.path.join(base_dir, "..", "temp")
    os.makedirs(temp_dir, exist_ok=True)

    # 确保输出目录存在
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # ========== 确定输入文件 ==========
    html_files = []
    input_base_dir = None

    if not os.path.exists(input_path):
        logging.error(f"输入路径不存在: {input_path}")
        return

    if os.path.isdir(input_path):
        # 处理目录输入
        input_base_dir = input_path
        try:
            html_files = [f for f in os.listdir(input_base_dir) if f.endswith('.html')]
            
            # 智能排序：提取文件名中的数字进行排序
            def extract_number(filename):
                match = re.search(r'file_(\d+)\.html', filename)
                return int(match.group(1)) if match else 0
            
            html_files.sort(key=extract_number)
            logging.info(f"在目录 '{input_base_dir}' 中找到并排序了 {len(html_files)} 个HTML文件")
        except Exception as e:
            logging.error(f"读取输入目录 '{input_base_dir}' 时出错: {e}")
            return
            
    elif os.path.isfile(input_path):
        # 处理单文件输入
        if input_path.endswith('.html'):
            input_base_dir = os.path.dirname(input_path)
            html_files = [os.path.basename(input_path)]
            logging.info(f"找到单个HTML文件: '{input_path}'")
        else:
            logging.error(f"输入文件不是HTML文件: '{input_path}'")
            return
    else:
        logging.error(f"输入路径不是有效的文件或目录: '{input_path}'")
        return
    
    logging.info(f"开始转换 {len(html_files)} 个文件到输出目录 '{output_dir}'，使用 {args.workers} 个工作线程")

    # 在工作线程间分配文件
    num_workers = min(args.workers, len(html_files))
    if num_workers == 0:
        logging.info("没有HTML文件需要处理")
        return
        
    # 为每个工作线程创建文件块
    file_chunks = [[] for _ in range(num_workers)]
    for i, html_file in enumerate(html_files):
        file_chunks[i % num_workers].append(html_file)

    tasks = [(chunk, input_base_dir, output_dir, temp_dir) for chunk in file_chunks]

    # 并行处理文件
    with concurrent.futures.ThreadPoolExecutor(max_workers=num_workers, thread_name_prefix='Converter') as executor:
        executor.map(process_files_worker, tasks)

    logging.info("所有文件处理完成")
    # 注意：临时文件保留用于调试，可根据需要取消注释以下代码进行清理
    # if os.path.exists(temp_dir):
    #     shutil.rmtree(temp_dir)
    #     logging.info("临时文件已清理")

if __name__ == "__main__":
    main()