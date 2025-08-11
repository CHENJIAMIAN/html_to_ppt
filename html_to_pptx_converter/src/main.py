import os
import time
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.by import By
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import logging
import argparse
import re

class SlideData:
    def __init__(self):
        self.background_image_path = None
        self.title_text = None
        self.title_geom = None
        self.subtitle_text = None
        self.subtitle_geom = None
        self.keyword_items = []

def init_driver():
    """
    Initializes the Selenium WebDriver.
    """
    options = webdriver.ChromeOptions()
    # 在调试时，可以注释掉 --headless 来观察浏览器行为
    options.add_argument("--headless")
    options.add_argument("--window-size=1280,720")
    options.add_argument("--hide-scrollbars")
    service = Service(ChromeDriverManager().install())
    driver = webdriver.Chrome(service=service, options=options)
    return driver

def extract_data_from_html(file_path):
    """
    Extracts structured data from each slide in the HTML file.
    """
    driver = init_driver()
    driver.get(f"file:///{os.path.abspath(file_path)}")
    time.sleep(1)

    slides_data = []
    slide_elements = driver.find_elements(By.CSS_SELECTOR, ".slide")

    for i, slide_element in enumerate(slide_elements):
        slide_data = SlideData()

        # ... [背景和文本提取逻辑保持不变] ...
        original_html = driver.execute_script("return arguments[0].innerHTML;", slide_element)
        html_without_text = original_html
        while re.search(r'>([^<>\s][^<>]*)<', html_without_text):
            html_without_text = re.sub(r'>([^<>]*)<', '><', html_without_text)
        driver.execute_script("arguments[0].innerHTML = arguments[1];", slide_element, html_without_text)
        time.sleep(0.2)
        try:
            screenshot_path = f"d:/Desktop/20250811/html_to_pptx_converter/temp/slide_{i}_bg.png"
            slide_element.screenshot(screenshot_path)
            slide_data.background_image_path = screenshot_path
        except Exception as e:
            logging.warning(f"Could not take background screenshot for slide {i}: {e}")
        driver.execute_script("arguments[0].innerHTML = arguments[1];", slide_element, original_html)
        time.sleep(0.1)
        # ... [背景和文本提取逻辑结束] ...

        # ... [Title 和 Subtitle 提取逻辑保持不变] ...
        try:
            title_element = slide_element.find_element(By.CSS_SELECTOR, ".title")
            slide_data.title_text = title_element.text
            slide_data.title_geom = {"x": title_element.location['x'], "y": title_element.location['y'], "width": title_element.size['width'], "height": title_element.size['height'], "font-size": title_element.value_of_css_property('font-size'), "color": title_element.value_of_css_property('color'), "font-weight": title_element.value_of_css_property('font-weight')}
        except Exception: pass
        try:
            subtitle_element = slide_element.find_element(By.CSS_SELECTOR, ".subtitle")
            slide_data.subtitle_text = subtitle_element.text
            slide_data.subtitle_geom = {"x": subtitle_element.location['x'], "y": subtitle_element.location['y'], "width": subtitle_element.size['width'], "height": subtitle_element.size['height'], "font-size": subtitle_element.value_of_css_property('font-size'), "color": subtitle_element.value_of_css_property('color')}
        except Exception: pass
        # ... [Title 和 Subtitle 提取逻辑结束] ...

        keyword_elements = slide_element.find_elements(By.CSS_SELECTOR, ".keyword-item")
        for j, keyword_element in enumerate(keyword_elements):
            keyword_item = {}
            try:
                icon_element = keyword_element.find_element(By.TAG_NAME, "i")

                # ##################################################################
                # --- 最终版：高分辨率 Icon 截图的可靠实现 ---
                # ##################################################################
                scale_factor = 5  # 放大5倍，获得非常高的分辨率

                # 1. 定义JS脚本，用于在页面左上角创建可渲染的大尺寸克隆体
                js_script = """
                    const targetElement = arguments[0];
                    const scale = arguments[1];

                    const style = window.getComputedStyle(targetElement);
                    const originalFontSizeStr = style.getPropertyValue('font-size');
                    const originalColor = style.getPropertyValue('color');
                    
                    if (!originalFontSizeStr) return null;
                    
                    const originalFontSize = parseFloat(originalFontSizeStr);

                    // 创建一个临时容器，并将其固定在屏幕左上角
                    const container = document.createElement('div');
                    container.id = 'temp-icon-container-for-screenshot';
                    container.style.position = 'absolute';
                    container.style.left = '0px';
                    container.style.top = '0px';
                    container.style.zIndex = '9999'; // 确保在最顶层
                    
                    const clone = targetElement.cloneNode(true);
                    
                    // 核心：放大字体大小，并继承颜色
                    clone.style.fontSize = (originalFontSize * scale) + 'px';
                    clone.style.color = originalColor;
                    
                    // 关键：确保背景透明，以便裁剪
                    clone.style.backgroundColor = 'transparent';
                    
                    container.appendChild(clone);
                    document.body.appendChild(container);
                    
                    return clone; // 返回这个被放大的克隆体
                """
                
                # 2. 执行脚本，获取到屏幕左上角的、已放大的克隆元素
                scaled_clone_element = driver.execute_script(js_script, icon_element, scale_factor)

                if scaled_clone_element:
                    # 关键：给浏览器一点时间来渲染新添加的元素
                    time.sleep(0.1)

                    try:
                        # 3. 对这个完美的克隆体进行截图
                        icon_path = f"d:/Desktop/20250811/html_to_pptx_converter/temp/slide_{i}_icon_{j}.png"
                        scaled_clone_element.screenshot(icon_path)

                        # 4. 使用Pillow精确裁剪掉截图周围的透明区域
                        with Image.open(icon_path) as img:
                            bbox = img.getbbox()
                            if bbox:
                                cropped_img = img.crop(bbox)
                                cropped_img.save(icon_path)
                        
                        keyword_item['icon_path'] = icon_path

                    finally:
                        # 5. 清理：无论成功与否，都移除临时容器
                        driver.execute_script(
                            "document.getElementById('temp-icon-container-for-screenshot').remove();"
                        )
                else:
                    logging.warning(f"Could not create a scalable clone for icon {j} on slide {i}.")

                # ##################################################################
                # --- 高分辨率截图逻辑结束 ---
                # ##################################################################

            except Exception as e:
                logging.warning(f"Could not find or process icon for keyword {j} on slide {i}: {e}")

            # ... [Keyword Title 和 Desc 提取逻辑保持不变] ...
            try:
                title_element = keyword_element.find_element(By.CSS_SELECTOR, ".keyword-title")
                keyword_item['title_text'] = title_element.text
                keyword_item['title_geom'] = {"x": title_element.location['x'], "y": title_element.location['y'], "width": title_element.size['width'], "height": title_element.size['height'], "font-size": title_element.value_of_css_property('font-size'), "color": title_element.value_of_css_property('color'), "font-weight": title_element.value_of_css_property('font-weight')}
            except Exception: pass
            try:
                desc_element = keyword_element.find_element(By.CSS_SELECTOR, ".keyword-desc")
                keyword_item['desc_text'] = desc_element.text
                keyword_item['desc_geom'] = {"x": desc_element.location['x'], "y": desc_element.location['y'], "width": desc_element.size['width'], "height": desc_element.size['height'], "font-size": desc_element.value_of_css_property('font-size'), "color": desc_element.value_of_css_property('color')}
            except Exception: pass
            # ... [Keyword Title 和 Desc 提取逻辑结束] ...
            
            if keyword_item:
                slide_data.keyword_items.append(keyword_item)
            
        slides_data.append(slide_data)
        
    driver.quit()
    return slides_data

# --- 以下函数保持不变 ---

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def px_to_emu(px):
    if isinstance(px, (int, float)):
        return int(px * 9525)
    return 0

def add_slide_with_background(prs, image_path):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        except Exception as e:
            logging.error(f"Failed to add background image {image_path}: {e}")
    return slide

def add_image(slide, image_path, x_emu, y_emu, width_emu):
    if image_path and os.path.exists(image_path):
        try:
            with Image.open(image_path) as img:
                aspect_ratio = img.height / img.width if img.width > 0 else 1
                height_emu = int(width_emu * aspect_ratio)
            slide.shapes.add_picture(image_path, x_emu, y_emu, width=width_emu, height=height_emu)
        except Exception as e:
             logging.error(f"Failed to add icon image {image_path}: {e}")

def add_textbox(slide, text, x_emu, y_emu, width_emu, height_emu, font_details):
    if not text: return
    try:
        textbox = slide.shapes.add_textbox(x_emu, y_emu, width_emu, height_emu)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text

        font = run.font
        if 'font-size' in font_details:
            try:
                 font_size_px = float(font_details['font-size'].replace("px", ""))
                 scale_factor = 0.95  # 在这里调整缩放比例，0.9代表90%，可根据需要改为0.85, 0.8等
                 font.size = Pt(int(font_size_px * scale_factor))
            except (ValueError, TypeError):
                 pass
        if 'color' in font_details:
            try:
                color_str = font_details['color'].replace("rgba(", "").replace(")", "").replace("rgb(", "")
                parts = [p.strip() for p in color_str.split(",")]
                r, g, b = int(parts[0]), int(parts[1]), int(parts[2])
                font.color.rgb = RGBColor(r, g, b)
            except Exception:
                pass
        if 'font-weight' in font_details and (str(font_details['font-weight']) == 'bold' or int(font_details['font-weight']) >= 700):
            font.bold = True
    except Exception as e:
        logging.error(f"Failed to add textbox with text '{text}': {e}")


def main():
    parser = argparse.ArgumentParser(description='Convert HTML to PPTX.')
    parser.add_argument('--input', type=str, required=True, help='Input HTML file path.')
    parser.add_argument('--output', type=str, required=True, help='Output PPTX file path.')
    args = parser.parse_args()

    logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
    logging.info(f"Starting conversion of {args.input} to {args.output}")

    temp_dir = "d:/Desktop/20250811/html_to_pptx_converter/temp"
    if not os.path.exists(temp_dir):
        os.makedirs(temp_dir)

    all_slides_data = extract_data_from_html(args.input)
    logging.info(f"Found {len(all_slides_data)} slides.")

    prs = create_presentation()

    for i, slide_data in enumerate(all_slides_data):
        logging.info(f"Processing slide {i+1}...")
        slide = add_slide_with_background(prs, slide_data.background_image_path)

        if slide_data.title_text and slide_data.title_geom:
            geom = slide_data.title_geom
            add_textbox(slide, slide_data.title_text, px_to_emu(geom['x']), px_to_emu(geom['y']), px_to_emu(geom['width']), px_to_emu(geom['height']), geom)

        if slide_data.subtitle_text and slide_data.subtitle_geom:
            geom = slide_data.subtitle_geom
            add_textbox(slide, slide_data.subtitle_text, px_to_emu(geom['x']), px_to_emu(geom['y']), px_to_emu(geom['width']), px_to_emu(geom['height']), geom)

        for keyword_item in slide_data.keyword_items:
            icon_width_px = 40
            icon_width_emu = px_to_emu(icon_width_px)
            icon_x, icon_y = 0, 0
            
            if keyword_item.get('title_geom'):
                title_geom = keyword_item['title_geom']
                icon_x = px_to_emu(title_geom['x']) - icon_width_emu - px_to_emu(10)
                title_center_y = title_geom['y'] + title_geom['height'] / 2
                with Image.open(keyword_item['icon_path']) as img:
                    aspect_ratio = img.height / img.width if img.width > 0 else 1
                    icon_height_px = icon_width_px * aspect_ratio
                icon_y = px_to_emu(title_center_y - icon_height_px / 2)
            
            if 'icon_path' in keyword_item:
                add_image(slide, keyword_item['icon_path'], icon_x, icon_y, icon_width_emu)
            
            if 'title_text' in keyword_item and 'title_geom' in keyword_item:
                geom = keyword_item['title_geom']
                add_textbox(slide, keyword_item['title_text'], px_to_emu(geom['x']), px_to_emu(geom['y']), px_to_emu(geom['width']), px_to_emu(geom['height']), geom)

            if 'desc_text' in keyword_item and 'desc_geom' in keyword_item:
                geom = keyword_item['desc_geom']
                add_textbox(slide, keyword_item['desc_text'], px_to_emu(geom['x']), px_to_emu(geom['y']), px_to_emu(geom['width']), px_to_emu(geom['height']), geom)

    output_dir = os.path.dirname(args.output)
    if not os.path.exists(output_dir) and output_dir:
        os.makedirs(output_dir)

    prs.save(args.output)
    logging.info(f"Presentation saved to {args.output}")

    try:
        for filename in os.listdir(temp_dir):
            file_path = os.path.join(temp_dir, filename)
            if os.path.isfile(file_path):
                os.remove(file_path)
        logging.info("Temporary files cleaned up.")
    except Exception as e:
        logging.warning(f"Could not clean up temp files: {e}")

if __name__ == "__main__":
    main()